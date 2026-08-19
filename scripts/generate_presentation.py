"""Build the Smart Surf School introduction decks from one source.

Produces ten files from ``presentation_content.py``:

    Surf_School_Tanitim.pptx / .pdf / .html          Turkish, screen
    Surf_School_Tanitim_Baski.pptx / .pdf            Turkish, print
    Surf_School_Intro_EN.pptx / .pdf / .html         English, screen
    Surf_School_Intro_EN_Print.pptx / .pdf           English, print

Why a generator instead of editing slides by hand
-------------------------------------------------
Four decks in two languages is the same content eight times over. Edited by
hand they drift apart within a week, and the PDF and HTML are derived from the
PPTX so any drift multiplies. Here the wording lives in one Python file, the
geometry lives in one renderer, and everything else is produced.

Screen and print are different documents, not the same one twice
----------------------------------------------------------------
The screen deck is dark: it is shown on a display, and a dark ground makes the
brand colour and the numbers carry. Printed, that same slide is a solid block of
ink — expensive, slow, and grey rather than black on most office printers. The
print variant is therefore rebuilt on white with darker text, not merely
"exported differently".

Requirements
------------
``python-pptx`` builds the PPTX. PDF and slide images come from PowerPoint via
COM automation, which is the only route on this machine — LibreOffice is not
installed. Without PowerPoint the PPTX files are still produced and the script
says clearly what it could not make.

Usage
-----
    python scripts/generate_presentation.py
    python scripts/generate_presentation.py --lang tr
    python scripts/generate_presentation.py --skip-pdf
"""

from __future__ import annotations

import argparse
import base64
import shutil
import sys
import tempfile
from pathlib import Path

BASE_DIR = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(BASE_DIR))
sys.path.insert(0, str(BASE_DIR / "scripts"))

for _stream in (sys.stdout, sys.stderr):
    if hasattr(_stream, "reconfigure"):
        _stream.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402
from presentation_content import (  # noqa: E402
    ACCENTS,
    BRAND,
    BRAND_DEEP,
    BRAND_LIGHT,
    DECKS,
    INK,
    INK_SOFT,
    OK,
    PAPER,
    PAPER_SOFT,
)

GREEN, RED, GREY, BOLD, RESET = "\033[92m", "\033[91m", "\033[90m", "\033[1m", "\033[0m"

# --- canvas -----------------------------------------------------------------
W = Inches(13.333)
H = Inches(7.5)
MARGIN = Inches(0.86)
CONTENT_W = W - 2 * MARGIN

FONT = "Segoe UI"
FONT_LIGHT = "Segoe UI Light"
FONT_MONO = "Consolas"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


class Theme:
    """Colours for one variant. Everything the renderer draws asks the theme."""

    def __init__(self, *, dark: bool):
        self.dark = dark
        if dark:
            self.bg_top = BRAND_DEEP
            self.bg_bottom = "041D30"
            self.panel = "0C3854"
            self.panel_edge = "12496B"
            self.title = "FFFFFF"
            self.body = "C7D9E6"
            self.muted = "7FA0B6"
            self.accent = BRAND_LIGHT
            self.rule = "17537A"
        else:
            self.bg_top = PAPER
            self.bg_bottom = PAPER
            self.panel = PAPER_SOFT
            self.panel_edge = "E2E8F0"
            self.title = INK
            self.body = INK_SOFT
            self.muted = "64748B"
            self.accent = BRAND
            self.rule = "CBD5E1"


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------
def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size=18,
    color="FFFFFF",
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line=1.25,
    space_after=0,
    caps=False,
    spacing=None,
):
    """One text box. Returns the shape so callers can tweak it further."""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0

    for index, chunk in enumerate(str(text).split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = line
        paragraph.space_after = Pt(space_after)
        run = paragraph.add_run()
        run.text = chunk.upper() if caps else chunk
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = rgb(color)
        if spacing is not None:
            # Letter-spacing is not in the python-pptx API; set it on the XML.
            run.font._rPr.set("spc", str(int(spacing * 100)))
    return box


def add_rect(slide, left, top, width, height, fill, *, edge=None, radius=0.04, shadow=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if edge:
        shape.line.color.rgb = rgb(edge)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius:
        try:
            shape.adjustments[0] = radius
        except (IndexError, KeyError):
            pass
    shape.shadow.inherit = shadow
    if shape.has_text_frame:
        shape.text_frame.text = ""
    return shape


def paint_background(slide, theme: Theme):
    """Full-bleed ground plus a thin brand rule along the top."""
    back = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    back.line.fill.background()
    back.shadow.inherit = False
    if theme.dark:
        fill = back.fill
        fill.gradient()
        fill.gradient_angle = 45.0
        stops = fill.gradient_stops
        stops[0].color.rgb = rgb(theme.bg_top)
        stops[0].position = 0.0
        stops[1].color.rgb = rgb(theme.bg_bottom)
        stops[1].position = 1.0
    else:
        back.fill.solid()
        back.fill.fore_color.rgb = rgb(theme.bg_top)
    add_rect(slide, 0, 0, W, Emu(int(Inches(0.055))), BRAND, radius=0)


def add_eyebrow(slide, theme: Theme, text: str, top=Inches(0.72)):
    add_text(
        slide, MARGIN, top, CONTENT_W, Inches(0.3), text,
        size=12, color=theme.accent, bold=True, caps=True, spacing=1.6,
    )


def add_title(slide, theme: Theme, text: str, top=Inches(1.06), size=38):
    lines = str(text).count("\n") + 1
    add_text(
        slide, MARGIN, top, CONTENT_W, Inches(0.72 * lines), text,
        size=size, color=theme.title, bold=True, line=1.1,
    )
    return top + Inches(0.62 * lines)


def add_lead(slide, theme: Theme, text: str, top, width=None):
    width = width or Emu(int(CONTENT_W * 0.86))
    add_text(slide, MARGIN, top, width, Inches(1.0), text, size=15, color=theme.body, line=1.45)
    lines = max(1, len(str(text)) // 105 + str(text).count("\n") + 1)
    return top + Inches(0.30 * lines) + Inches(0.14)


def add_note(slide, theme: Theme, text: str):
    """Footnote pinned just above the footer."""
    add_rect(slide, MARGIN, Inches(6.36), Inches(0.045), Inches(0.44), theme.accent, radius=0)
    add_text(
        slide, MARGIN + Inches(0.22), Inches(6.34), Emu(int(CONTENT_W * 0.9)), Inches(0.6),
        text, size=11.5, color=theme.muted, line=1.35,
    )


def add_footer(slide, theme: Theme, index: int, total: int, label: str):
    add_rect(slide, MARGIN, Inches(6.97), CONTENT_W, Emu(9525), theme.rule, radius=0)
    add_text(
        slide, MARGIN, Inches(7.06), Inches(8.0), Inches(0.3), label,
        size=10, color=theme.muted,
    )
    add_text(
        slide, W - MARGIN - Inches(2.0), Inches(7.06), Inches(2.0), Inches(0.3),
        f"{index} / {total}", size=10, color=theme.muted, align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Slide kinds
# ---------------------------------------------------------------------------
def draw_title(slide, theme: Theme, data: dict):
    band_w = Inches(4.6)
    band_left = W - band_w
    if theme.dark:
        # A slightly lighter band on the right gives the title slide depth
        # without an image that would bloat every export.
        add_rect(slide, band_left, Inches(0.055), band_w, H, "0A3350", radius=0)

    # The band carries the headline figures rather than sitting empty.
    for offset, (value, caption) in enumerate(data.get("highlights", [])):
        top = Inches(2.05) + Inches(1.32) * offset
        accent = ACCENTS[offset % len(ACCENTS)] if theme.dark else BRAND
        add_rect(slide, band_left + Inches(0.9), top, Inches(0.05), Inches(0.86), accent, radius=0)
        add_text(
            slide, band_left + Inches(1.16), top - Inches(0.06), Inches(3.0), Inches(0.6),
            value, size=30, color=theme.title, bold=True, line=1.0,
        )
        add_text(
            slide, band_left + Inches(1.16), top + Inches(0.5), Inches(3.0), Inches(0.4),
            caption, size=11.5, color=theme.muted, line=1.2,
        )

    add_rect(slide, MARGIN, Inches(2.05), Inches(0.9), Inches(0.07), theme.accent, radius=0)
    add_text(
        slide, MARGIN, Inches(1.55), Inches(9.0), Inches(0.34), data["eyebrow"],
        size=12.5, color=theme.accent, bold=True, caps=True, spacing=2.0,
    )
    add_text(
        slide, MARGIN, Inches(2.45), Inches(9.4), Inches(2.1), data["title"],
        size=52, color=theme.title, bold=True, line=1.06,
    )
    add_text(
        slide, MARGIN, Inches(4.72), Inches(8.6), Inches(1.1), data["subtitle"],
        size=16, color=theme.body, line=1.5,
    )
    add_text(
        slide, MARGIN, Inches(6.2), Inches(9.0), Inches(0.4), data["meta"],
        size=12, color=theme.muted,
    )


def draw_metrics(slide, theme: Theme, data: dict):
    add_eyebrow(slide, theme, data["eyebrow"])
    top = add_title(slide, theme, data["title"])
    if data.get("lead"):
        top = add_lead(slide, theme, data["lead"], top)

    items = data["metrics"]
    columns = 4
    rows = (len(items) + columns - 1) // columns
    gap = Inches(0.26)
    tile_w = Emu(int((CONTENT_W - gap * (columns - 1)) / columns))
    start = top + Inches(0.18)

    # Grow into the free space, but only so far: a tile taller than its content
    # is just a big empty box. Whatever height it lands on, the value and its
    # caption sit as one block in the middle of it.
    bottom = Inches(6.2) if data.get("note") else Inches(6.74)
    tile_h = min(
        Emu(int((bottom - start - gap * (rows - 1)) / rows)),
        Inches(2.5) if rows == 1 else Inches(2.0),
    )

    for index, (value, caption) in enumerate(items):
        row, column = divmod(index, columns)
        left = MARGIN + Emu(int((tile_w + gap) * column))
        tile_top = start + Emu(int((tile_h + gap) * row))
        add_rect(slide, left, tile_top, tile_w, tile_h, theme.panel, edge=theme.panel_edge)
        accent = ACCENTS[index % len(ACCENTS)] if theme.dark else BRAND
        add_rect(slide, left, tile_top, tile_w, Inches(0.055), accent, radius=0)
        value_h = Inches(0.72 if rows > 1 else 0.86)
        caption_lines = caption.count("\n") + 1
        caption_h = Inches(0.26) * caption_lines + Inches(0.06)
        block_top = tile_top + Emu(int((tile_h - value_h - caption_h) / 2)) + Inches(0.04)

        add_text(
            slide, left + Inches(0.24), block_top, tile_w - Inches(0.48), value_h,
            value, size=33 if rows > 1 else 42,
            color=theme.title, bold=True, align=PP_ALIGN.LEFT, line=1.0,
        )
        add_text(
            slide, left + Inches(0.24), block_top + value_h, tile_w - Inches(0.48),
            caption_h, caption,
            size=11.5 if rows > 1 else 13, color=theme.muted, line=1.3,
        )

    if data.get("note"):
        add_note(slide, theme, data["note"])


def draw_cards(slide, theme: Theme, data: dict):
    add_eyebrow(slide, theme, data["eyebrow"])
    add_title(slide, theme, data["title"])

    cards = data["cards"]
    columns = 3 if len(cards) > 4 else 2
    rows = (len(cards) + columns - 1) // columns
    gap = Inches(0.26)
    card_w = Emu(int((CONTENT_W - gap * (columns - 1)) / columns))
    start = Inches(2.16)

    # Fill the space down to the note (or the footer) instead of leaving a band
    # of dead ground under the last row.
    bottom = Inches(6.2) if data.get("note") else Inches(6.74)
    card_h = Emu(int((bottom - start - gap * (rows - 1)) / rows))

    for index, (heading, body) in enumerate(cards):
        row, column = divmod(index, columns)
        left = MARGIN + Emu(int((card_w + gap) * column))
        top = start + Emu(int((card_h + gap) * row))
        add_rect(slide, left, top, card_w, card_h, theme.panel, edge=theme.panel_edge)
        accent = ACCENTS[index % len(ACCENTS)] if theme.dark else BRAND
        add_rect(slide, left, top, Inches(0.055), card_h, accent, radius=0)
        add_text(
            slide, left + Inches(0.28), top + Inches(0.22), card_w - Inches(0.5),
            Inches(0.4), heading, size=15.5, color=theme.title, bold=True, line=1.15,
        )
        add_text(
            slide, left + Inches(0.28), top + Inches(0.66), card_w - Inches(0.5),
            card_h - Inches(0.86), body, size=11.5, color=theme.body, line=1.38,
        )

    if data.get("note"):
        add_note(slide, theme, data["note"])


def draw_bullets(slide, theme: Theme, data: dict):
    add_eyebrow(slide, theme, data["eyebrow"])
    top = add_title(slide, theme, data["title"])
    if data.get("lead"):
        top = add_lead(slide, theme, data["lead"], top, width=Emu(int(CONTENT_W * 0.82)))

    items = data["items"]
    row_h = Inches(0.63)
    start = max(top + Inches(0.1), Inches(2.62))
    available = Inches(6.2) - start
    if row_h * len(items) > available:
        row_h = Emu(int(available / len(items)))

    for index, (heading, body) in enumerate(items):
        top_i = start + Emu(int(row_h * index))
        accent = ACCENTS[index % len(ACCENTS)] if theme.dark else BRAND
        add_rect(slide, MARGIN, top_i + Inches(0.08), Inches(0.2), Inches(0.2), accent, radius=0.5)
        add_text(
            slide, MARGIN + Inches(0.42), top_i, Inches(3.0), Inches(0.36), heading,
            size=13.5, color=theme.title, bold=True,
        )
        add_text(
            slide, MARGIN + Inches(3.55), top_i, CONTENT_W - Inches(3.55), Inches(0.5),
            body, size=12, color=theme.body, line=1.3,
        )

    if data.get("note"):
        add_note(slide, theme, data["note"])


def draw_split(slide, theme: Theme, data: dict):
    add_eyebrow(slide, theme, data["eyebrow"])
    add_title(slide, theme, data["title"])

    left_w = Inches(5.05)
    add_text(
        slide, MARGIN, Inches(2.26), left_w, Inches(3.4), data["lead"],
        size=14, color=theme.body, line=1.5,
    )

    rows = data["rows"]
    panel_left = MARGIN + left_w + Inches(0.6)
    panel_w = W - MARGIN - panel_left
    panel_top = Inches(2.16)

    # The panel must finish above the note and the footer rule. Row height is
    # derived from what is actually left rather than assumed, so a nine-row
    # table cannot run off the bottom of the slide.
    panel_bottom = Inches(6.18) if data.get("note") else Inches(6.72)
    padding = Inches(0.2)
    available = panel_bottom - panel_top - 2 * padding
    row_h = min(Inches(0.52), Emu(int(available / max(1, len(rows)))))
    total_h = Emu(int(row_h * len(rows))) + 2 * padding
    add_rect(slide, panel_left, panel_top, panel_w, total_h, theme.panel, edge=theme.panel_edge)

    # Values are right-aligned but need room: several of them are two words that
    # would otherwise wrap onto a cramped second line.
    value_w = Inches(2.35)
    label_w = panel_w - value_w - Inches(0.72)
    compact = row_h < Inches(0.46)

    for index, (label, value) in enumerate(rows):
        top = panel_top + padding + Emu(int(row_h * index))
        add_text(
            slide, panel_left + Inches(0.28), top, label_w, row_h,
            label, size=11.5 if compact else 12.5, color=theme.title, bold=True,
            anchor=MSO_ANCHOR.MIDDLE, line=1.1,
        )
        add_text(
            slide, panel_left + panel_w - value_w - Inches(0.28), top, value_w, row_h,
            value, size=10.5 if compact else 11.5, color=theme.muted,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, line=1.1,
        )
        if index < len(rows) - 1:
            add_rect(
                slide, panel_left + Inches(0.28), top + row_h - Emu(9525),
                panel_w - Inches(0.56), Emu(9525), theme.rule, radius=0,
            )

    if data.get("note"):
        add_note(slide, theme, data["note"])


def draw_screens(slide, theme: Theme, data: dict):
    """One or two real screenshots of the application, with captions.

    The images are PNGs captured from the running app (see
    ``scripts/capture_screenshots.py``); paths in the content file are
    relative to the repository root. Screenshots are 16:9, the same ratio as
    the slide, so width and height are forced together without distortion.
    """
    add_eyebrow(slide, theme, data["eyebrow"])
    top = add_title(slide, theme, data["title"])
    if data.get("lead"):
        top = add_lead(slide, theme, data["lead"], top)

    shots = data["shots"]
    bottom = Inches(6.2) if data.get("note") else Inches(6.74)
    caption_h = Inches(0.46)
    start = top + Inches(0.12)
    avail_h = bottom - start - caption_h
    aspect = 16 / 9
    gap = Inches(0.3)

    if len(shots) == 1:
        img_h = min(int(avail_h), int(CONTENT_W / aspect))
        img_w = int(img_h * aspect)
        columns = [(MARGIN + int((CONTENT_W - img_w) / 2), img_w)]
    else:
        img_w = int((CONTENT_W - gap * (len(shots) - 1)) / len(shots))
        img_h = int(img_w / aspect)
        if img_h > int(avail_h):
            img_h = int(avail_h)
            img_w = int(img_h * aspect)
        columns = [
            (MARGIN + int((CONTENT_W - (img_w * len(shots) + int(gap) * (len(shots) - 1))) / 2)
             + (img_w + int(gap)) * index, img_w)
            for index in range(len(shots))
        ]

    # Centre the block vertically in the free band rather than leaving all the
    # slack under the captions.
    slack = int(avail_h) - img_h
    if slack > 0:
        start += int(slack / 2)

    for (left, width), (path, caption) in zip(columns, shots):
        source = BASE_DIR / path
        if not source.exists():
            raise FileNotFoundError(
                f"Screenshot missing: {source}\n"
                "Run scripts/capture_screenshots.py against a seeded dev server first."
            )
        picture = slide.shapes.add_picture(str(source), left, start, width, img_h)
        picture.line.color.rgb = rgb(theme.panel_edge)
        picture.line.width = Pt(1.25)
        picture.shadow.inherit = False
        add_text(
            slide, left, start + img_h + Inches(0.1), width, caption_h, caption,
            size=11.5, color=theme.muted, line=1.3, align=PP_ALIGN.CENTER,
        )

    if data.get("note"):
        add_note(slide, theme, data["note"])


def draw_quote(slide, theme: Theme, data: dict):
    add_eyebrow(slide, theme, data["eyebrow"])
    add_rect(slide, MARGIN, Inches(1.55), Inches(0.075), Inches(3.4), OK if theme.dark else BRAND, radius=0)
    add_text(
        slide, MARGIN + Inches(0.42), Inches(1.6), Inches(10.6), Inches(1.5),
        data["quote"], size=31, color=theme.title, bold=True, line=1.22,
    )
    add_text(
        slide, MARGIN + Inches(0.42), Inches(3.55), Inches(9.6), Inches(2.4),
        data["body"], size=15, color=theme.body, line=1.55,
    )


def draw_closing(slide, theme: Theme, data: dict):
    add_eyebrow(slide, theme, data["eyebrow"])

    panel_left = Inches(7.1)
    # Keep the headline clear of the panel: at 40pt a long Turkish line ran
    # underneath it and the last word disappeared.
    title_w = panel_left - MARGIN - Inches(0.45)
    add_text(
        slide, MARGIN, Inches(1.5), title_w, Inches(2.1), data["title"],
        size=36, color=theme.title, bold=True, line=1.12,
    )

    panel_w = W - MARGIN - panel_left
    points = data["points"]
    row_h = Inches(0.72)
    panel_top = Inches(1.5)
    add_rect(
        slide, panel_left, panel_top, panel_w,
        Emu(int(row_h * len(points))) + Inches(0.4), theme.panel, edge=theme.panel_edge,
    )
    for index, point in enumerate(points):
        top = panel_top + Inches(0.2) + Emu(int(row_h * index))
        accent = ACCENTS[index % len(ACCENTS)] if theme.dark else BRAND
        add_rect(slide, panel_left + Inches(0.26), top + Inches(0.16), Inches(0.18), Inches(0.18), accent, radius=0.5)
        add_text(
            slide, panel_left + Inches(0.62), top, panel_w - Inches(0.9), Inches(0.6),
            point, size=12.5, color=theme.body, line=1.3, anchor=MSO_ANCHOR.TOP,
        )

    add_rect(slide, MARGIN, Inches(5.5), Inches(0.9), Inches(0.06), theme.accent, radius=0)
    add_text(
        slide, MARGIN, Inches(5.86), Inches(6.2), Inches(0.9), data["meta"],
        size=12.5, color=theme.muted, line=1.5,
    )


DRAW = {
    "title": draw_title,
    "metrics": draw_metrics,
    "cards": draw_cards,
    "bullets": draw_bullets,
    "split": draw_split,
    "screens": draw_screens,
    "quote": draw_quote,
    "closing": draw_closing,
}


# ---------------------------------------------------------------------------
# Deck assembly
# ---------------------------------------------------------------------------
def build_pptx(slides: list[dict], theme: Theme, footer_label: str, out_path: Path) -> Path:
    presentation = Presentation()
    presentation.slide_width = W
    presentation.slide_height = H
    blank = presentation.slide_layouts[6]

    total = len(slides)
    for index, data in enumerate(slides, start=1):
        slide = presentation.slides.add_slide(blank)
        paint_background(slide, theme)
        DRAW[data["kind"]](slide, theme, data)
        if data["kind"] not in {"title", "closing"}:
            add_footer(slide, theme, index, total, footer_label)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(out_path))
    return out_path


# ---------------------------------------------------------------------------
# PowerPoint export
# ---------------------------------------------------------------------------
PPT_PDF = 32
PPT_PNG = 18


def _powerpoint():
    import win32com.client

    app = win32com.client.Dispatch("PowerPoint.Application")
    return app


def export_pdf(pptx_path: Path, pdf_path: Path) -> bool:
    """Export a PDF through PowerPoint. Returns False when it is unavailable."""
    try:
        app = _powerpoint()
    except Exception as exc:  # noqa: BLE001 - PowerPoint is optional tooling
        print(f"  {RED}PowerPoint unavailable{RESET} ({type(exc).__name__}) — PDF skipped")
        return False

    deck = None
    try:
        deck = app.Presentations.Open(str(pptx_path), WithWindow=False)
        deck.SaveAs(str(pdf_path), PPT_PDF)
        return True
    except Exception as exc:  # noqa: BLE001
        print(f"  {RED}PDF export failed{RESET}: {exc}")
        return False
    finally:
        try:
            if deck is not None:
                deck.Close()
        except Exception:  # noqa: BLE001, S110 - deliberate best-effort cleanup; a failure here must not break the caller
            pass


def export_pngs(pptx_path: Path, target_dir: Path, width=1920, height=1080) -> list[Path]:
    """Render every slide to PNG. Empty list when PowerPoint is unavailable."""
    try:
        app = _powerpoint()
    except Exception:  # noqa: BLE001
        return []

    target_dir.mkdir(parents=True, exist_ok=True)
    deck = None
    try:
        deck = app.Presentations.Open(str(pptx_path), WithWindow=False)
        deck.SaveAs(str(target_dir), PPT_PNG)
    except Exception as exc:  # noqa: BLE001
        print(f"  {RED}PNG export failed{RESET}: {exc}")
        return []
    finally:
        try:
            if deck is not None:
                deck.Close()
        except Exception:  # noqa: BLE001, S110 - deliberate best-effort cleanup; a failure here must not break the caller
            pass

    def order(path: Path) -> int:
        digits = "".join(ch for ch in path.stem if ch.isdigit())
        return int(digits) if digits else 0

    return sorted(target_dir.glob("*.PNG"), key=order) or sorted(
        target_dir.glob("*.png"), key=order
    )


# ---------------------------------------------------------------------------
# HTML deck
# ---------------------------------------------------------------------------
HTML_SHELL = """<!doctype html>
<html lang="{lang}">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1, viewport-fit=cover">
<title>{title}</title>
<style>
  :root {{ --bg:#04121f; --panel:#0c2a3f; --ink:#eef6fb; --muted:#7fa0b6; --brand:#2cbeff; }}
  * {{ box-sizing:border-box; -webkit-tap-highlight-color:transparent; }}
  html,body {{ margin:0; padding:0; background:var(--bg); color:var(--ink);
    font-family:"Segoe UI",system-ui,-apple-system,sans-serif; overscroll-behavior:none; }}
  header {{ position:fixed; inset:0 0 auto 0; height:52px; display:flex; align-items:center;
    gap:12px; padding:0 16px; background:rgba(4,18,31,.92); backdrop-filter:blur(8px);
    border-bottom:1px solid #10374f; z-index:10; }}
  header .dot {{ width:26px; height:26px; border-radius:8px; background:var(--brand);
    display:grid; place-items:center; flex:none; }}
  header h1 {{ font-size:14px; font-weight:600; margin:0; white-space:nowrap; overflow:hidden;
    text-overflow:ellipsis; }}
  header .count {{ margin-left:auto; font-size:12px; color:var(--muted);
    font-variant-numeric:tabular-nums; flex:none; }}
  main {{ padding:52px 0 76px; }}
  .stage {{ display:flex; overflow-x:auto; scroll-snap-type:x mandatory;
    scrollbar-width:none; }}
  .stage::-webkit-scrollbar {{ display:none; }}
  .slide {{ flex:0 0 100%; scroll-snap-align:center; padding:14px; }}
  .slide img {{ width:100%; height:auto; display:block; border-radius:12px;
    box-shadow:0 10px 34px rgba(0,0,0,.5); background:#08243a; }}
  nav {{ position:fixed; inset:auto 0 0 0; height:64px; display:flex; align-items:center;
    justify-content:center; gap:10px; padding:0 16px;
    background:rgba(4,18,31,.94); border-top:1px solid #10374f; z-index:10; }}
  nav button {{ appearance:none; border:1px solid #17537a; background:var(--panel);
    color:var(--ink); font:inherit; font-size:13px; padding:9px 16px; border-radius:9px;
    cursor:pointer; min-width:104px; }}
  nav button:disabled {{ opacity:.36; cursor:default; }}
  nav button:active {{ transform:translateY(1px); }}
  .rail {{ display:flex; gap:5px; align-items:center; max-width:38vw; overflow:hidden; }}
  .pip {{ width:6px; height:6px; border-radius:50%; background:#1d5c81; flex:none;
    transition:background .18s,transform .18s; }}
  .pip.on {{ background:var(--brand); transform:scale(1.5); }}
  .hint {{ text-align:center; color:var(--muted); font-size:11.5px; padding:2px 0 12px; }}
  @media (min-width:900px) {{
    .slide {{ padding:22px 40px; }}
    header h1 {{ font-size:15px; }}
  }}
  @media print {{
    header,nav,.hint {{ display:none; }}
    main {{ padding:0; }}
    .stage {{ display:block; overflow:visible; }}
    .slide {{ padding:0; page-break-after:always; }}
    .slide img {{ border-radius:0; box-shadow:none; }}
  }}
</style>
</head>
<body>
<header>
  <span class="dot">
    <svg width="15" height="15" viewBox="0 0 24 24" fill="none" stroke="#04121f"
         stroke-width="2.4" stroke-linecap="round" stroke-linejoin="round">
      <path d="M2 12c2 0 2 2.2 4 2.2S8 12 10 12s2 2.2 4 2.2S16 12 18 12s2 2.2 4 2.2"/>
      <path d="M2 18c2 0 2 2.2 4 2.2S8 18 10 18s2 2.2 4 2.2S16 18 18 18s2 2.2 4 2.2"/>
    </svg>
  </span>
  <h1>{title}</h1>
  <span class="count"><b id="now">1</b>{of}{total}</span>
</header>

<main>
  <div class="stage" id="stage">
{slides}
  </div>
  <p class="hint">{hint}</p>
</main>

<nav>
  <button id="prev" type="button">&#8592; {prev}</button>
  <div class="rail" id="rail"></div>
  <button id="next" type="button">{next} &#8594;</button>
</nav>

<script>
(function () {{
  var stage = document.getElementById('stage');
  var slides = Array.prototype.slice.call(stage.children);
  var rail = document.getElementById('rail');
  var now = document.getElementById('now');
  var prev = document.getElementById('prev');
  var next = document.getElementById('next');
  var index = 0;

  slides.forEach(function (_, i) {{
    var pip = document.createElement('span');
    pip.className = 'pip' + (i === 0 ? ' on' : '');
    pip.addEventListener('click', function () {{ go(i); }});
    rail.appendChild(pip);
  }});
  var pips = Array.prototype.slice.call(rail.children);

  function paint() {{
    now.textContent = index + 1;
    pips.forEach(function (p, i) {{ p.classList.toggle('on', i === index); }});
    prev.disabled = index === 0;
    next.disabled = index === slides.length - 1;
    var active = pips[index];
    if (active) rail.scrollLeft = active.offsetLeft - rail.clientWidth / 2;
  }}

  function go(i) {{
    index = Math.max(0, Math.min(slides.length - 1, i));
    stage.scrollTo({{ left: slides[index].offsetLeft, behavior: 'smooth' }});
    paint();
  }}

  prev.addEventListener('click', function () {{ go(index - 1); }});
  next.addEventListener('click', function () {{ go(index + 1); }});

  document.addEventListener('keydown', function (e) {{
    if (e.key === 'ArrowRight' || e.key === 'PageDown' || e.key === ' ') {{ e.preventDefault(); go(index + 1); }}
    if (e.key === 'ArrowLeft' || e.key === 'PageUp') {{ e.preventDefault(); go(index - 1); }}
    if (e.key === 'Home') go(0);
    if (e.key === 'End') go(slides.length - 1);
  }});

  var timer;
  stage.addEventListener('scroll', function () {{
    clearTimeout(timer);
    timer = setTimeout(function () {{
      var i = Math.round(stage.scrollLeft / stage.clientWidth);
      if (i !== index) {{ index = i; paint(); }}
    }}, 90);
  }});

  paint();
}})();
</script>
</body>
</html>
"""


def build_html(pngs: list[Path], deck: dict, lang: str, out_path: Path) -> Path:
    """Self-contained HTML deck with every slide embedded as base64."""
    parts = []
    for index, png in enumerate(pngs, start=1):
        encoded = base64.b64encode(png.read_bytes()).decode("ascii")
        parts.append(
            f'    <figure class="slide">'
            f'<img src="data:image/png;base64,{encoded}" '
            f'alt="{deck["html_title"]} — {index}/{len(pngs)}" loading="lazy"></figure>'
        )

    html = HTML_SHELL.format(
        lang=lang,
        title=deck["html_title"],
        slides="\n".join(parts),
        total=len(pngs),
        of=deck["of"],
        prev=deck["nav_prev"],
        next=deck["nav_next"],
        hint=deck["hint"],
    )
    out_path.write_text(html, encoding="utf-8")
    return out_path


# ---------------------------------------------------------------------------
# Orchestration
# ---------------------------------------------------------------------------
def produce(lang: str, out_dir: Path, *, skip_pdf: bool, skip_html: bool) -> list[Path]:
    deck = DECKS[lang]
    slides = deck["slides"]
    stem = deck["stem"]
    made: list[Path] = []

    label = deck["html_title"]

    print(f"\n{BOLD}{lang.upper()} — {len(slides)} slides{RESET}")

    # --- screen -----------------------------------------------------------
    screen_pptx = build_pptx(slides, Theme(dark=True), label, out_dir / f"{stem}.pptx")
    made.append(screen_pptx)
    print(f"  {GREEN}ok{RESET}  {screen_pptx.name}  {GREY}{screen_pptx.stat().st_size:,} bytes{RESET}")

    # --- print ------------------------------------------------------------
    print_pptx = build_pptx(
        slides, Theme(dark=False), label, out_dir / f"{stem}{deck['print_suffix']}.pptx"
    )
    made.append(print_pptx)
    print(f"  {GREEN}ok{RESET}  {print_pptx.name}  {GREY}{print_pptx.stat().st_size:,} bytes{RESET}")

    if skip_pdf:
        return made

    for source in (screen_pptx, print_pptx):
        pdf = source.with_suffix(".pdf")
        if export_pdf(source, pdf) and pdf.exists():
            made.append(pdf)
            print(f"  {GREEN}ok{RESET}  {pdf.name}  {GREY}{pdf.stat().st_size:,} bytes{RESET}")

    if skip_html:
        return made

    temp = Path(tempfile.mkdtemp(prefix=f"surf_slides_{lang}_"))
    try:
        pngs = export_pngs(screen_pptx, temp)
        if pngs:
            html = build_html(pngs, deck, lang, out_dir / f"{stem}.html")
            made.append(html)
            print(
                f"  {GREEN}ok{RESET}  {html.name}  "
                f"{GREY}{html.stat().st_size:,} bytes · {len(pngs)} slides embedded{RESET}"
            )
        else:
            print(f"  {RED}--{RESET}  {stem}.html  {GREY}no slide images available{RESET}")
    finally:
        shutil.rmtree(temp, ignore_errors=True)

    return made


def main() -> int:
    parser = argparse.ArgumentParser(description="Build the introduction decks.")
    parser.add_argument("--lang", choices=["tr", "en", "all"], default="all")
    parser.add_argument("--out", default=str(BASE_DIR / "docs" / "presentation"))
    parser.add_argument("--skip-pdf", action="store_true")
    parser.add_argument("--skip-html", action="store_true")
    options = parser.parse_args()

    out_dir = Path(options.out)
    out_dir.mkdir(parents=True, exist_ok=True)

    print(f"{BOLD}Smart Surf School — introduction decks{RESET}")
    print(f"{GREY}Output: {out_dir}{RESET}")

    languages = ["tr", "en"] if options.lang == "all" else [options.lang]
    produced: list[Path] = []
    for lang in languages:
        produced += produce(
            lang, out_dir, skip_pdf=options.skip_pdf, skip_html=options.skip_html
        )

    total = sum(p.stat().st_size for p in produced if p.exists())
    print(f"\n{BOLD}{'=' * 62}{RESET}")
    print(f"  {GREEN}{len(produced)} file(s){RESET}  {GREY}{total:,} bytes total{RESET}")
    print(f"{BOLD}{'=' * 62}{RESET}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
